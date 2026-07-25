#!/usr/bin/env python3
"""Plate-model slides (16:9) for the 03-07-2026 Emilio deck — analytical-model section.

Reuses the deck's house style (AE method: assertion headline + evidence, all text
black, Calibri, one succinct caption per plot stating what is shown + quantity/units).
Plots carry NO in-canvas text annotation (scientific-figures §2) — the quantitative /
regime information lives in the captions here.

Run: /local/data/public/skcb2/tddft/venv/bin/python3 build_plate_model_slides.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(HERE, "figs")
EQ = os.path.join(HERE, "eqns")
OUT = os.path.join(HERE, "plate_model_slides.pptx")

SW, SH, MARGIN = 10.0, 5.625, 0.08
FONT = "Calibri"; BLACK = RGBColor(0, 0, 0)
prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def _tb(slide, x, y, w, h, text, size, *, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = BLACK
    return box

def headline(slide, text):
    _tb(slide, MARGIN, 0.06, SW - 2 * MARGIN, 1.0, text, 22, bold=True)

def caption(slide, x, y, w, text):
    _tb(slide, x, y, w, 0.6, text, 11.5, align=PP_ALIGN.CENTER)

def _imgsize(path):
    from PIL import Image
    with Image.open(path) as im:
        return im.size

def image(slide, path, x, y, w, h, cap=None):
    cap_h = 0.55 if cap else 0.0
    box_h = h - cap_h
    iw, ih = _imgsize(path); ar = iw / ih
    if w / box_h > ar:
        dw = box_h * ar; dh = box_h
    else:
        dw = w; dh = w / ar
    dx = x + (w - dw) / 2; dy = y + (box_h - dh) / 2
    slide.shapes.add_picture(path, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
    if cap:
        caption(slide, x, y + box_h + 0.02, w, cap)

def new_slide(title):
    s = prs.slides.add_slide(BLANK); headline(s, title); return s

# --------------------------------------------------------------- slide 1: model
s = new_slide("A 1-D plate model of the slab predicts the projectile–slab energy")
eqs = ["distance.png", "U_total.png", "phi.png", "identity.png"]
y = 1.30
for e in eqs:
    p = os.path.join(EQ, e)
    iw, ih = _imgsize(p); ar = iw / ih
    dh = 0.72; dw = dh * ar
    if dw > 8.4: dw = 8.4; dh = dw / ar
    s.shapes.add_picture(p, Inches((SW - dw) / 2), Inches(y), Inches(dw), Inches(dh))
    y += dh + 0.14
caption(s, 1.0, y + 0.02, 8.0,
        "φ(z) from planar-averaged charge ρ = n₊ − n_e; point vs Gaussian-packet "
        "energies differ only by −2πqs²ρ; image term added for r ≳ 3 Bohr.")

# --------------------------------------------------------------- slide 2: model vs slab
s = new_slide("The model reproduces the DFT slab: Friedel peaks, spill-out, surface dipole")
image(s, os.path.join(FIGS, "plate_model_density_potential.png"), 2.7, 1.20, 4.6, 4.25,
      cap="Top: planar-averaged electron density n_e(z) (DFT) and jellium background "
          "n₊(z) [10⁻³ a₀⁻³]. Bottom: electrostatic potential φ(z) [eV]. Shaded = slab "
          "|z|<12.5 Bohr; dashed = Friedel first maxima (|z|≈7.9 Bohr). Symmetric dipole "
          "barrier ≈1.8 eV.")

# --------------------------------------------------------------- slide 3: wp weighting
s = new_slide("Wavepacket vs point charge differ only by −2πqs²ρ — meV, only at the surface")
image(s, os.path.join(FIGS, "plate_model_wp_weighting.png"), 1.9, 1.20, 6.2, 4.25,
      cap="U_wp − U_pt vs packet centre z_c [meV]: numeric (solid) vs the analytic "
          "−2πqs²ρ(z) identity (dashed). They agree to sub-meV; the difference ∝ local "
          "net charge ρ, so it vanishes in the neutral interior and vacuum, peaking "
          "(≈−9 meV) at the surface. σ = 0.5 Bohr, s = σ/√2.")

# --------------------------------------------------------------- slide 4: U(r)
s = new_slide("In vacuum the attraction is image physics (∝1/r); inside, the surface barrier caps it")
image(s, os.path.join(FIGS, "plate_model_U_of_r.png"), 1.9, 1.20, 6.2, 4.25,
      cap="Interaction energy U(r) [eV] vs distance r=|z_c|−a from the slab face: static "
          "point qφ, static wavepacket q(g_s∗φ) (indistinguishable, <16 meV apart), image "
          "−q²/4(r−z_im), and total. Vacuum r≳7: image-dominated 1/r; interior r<0: dipole "
          "barrier plateau. WP internal zero-point (81.6 eV) excluded.")

prs.save(OUT)
print("wrote", OUT, "—", len(prs.slides._sldIdLst), "slides")
